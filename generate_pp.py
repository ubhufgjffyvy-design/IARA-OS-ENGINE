from pptx import Presentation
from pptx.util import Inches

# Criar uma apresentação
prs = Presentation()

# Slide 1: Título
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Apresentação sobre a IA que cria aplicativos e arquivos sem erros"
subtitle.text = "Uma visão geral das capacidades da IA IARA-OS-ENGINE"

# Slide 2: Introdução
bullet_slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = 'Introdução'
tf = body_shape.text_frame
tf.text = 'A IA IARA-OS-ENGINE é uma assistente avançada de programação que ajuda no desenvolvimento de software.'
p = tf.add_paragraph()
p.text = 'Ela é capaz de gerar código, criar arquivos e corrigir erros de forma eficiente.'
p.level = 0

# Slide 3: Funcionalidades
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = 'Funcionalidades'
tf = body_shape.text_frame
tf.text = 'A IA faz aplicativos completos sem erros.'
p = tf.add_paragraph()
p.text = 'Ela cria arquivos de código fonte, configurações e documentação.'
p.level = 0
p = tf.add_paragraph()
p.text = 'Também corrige arquivos existentes, identificando e resolvendo bugs.'
p.level = 0

# Slide 4: Capacidades de Conversação
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = 'Capacidades de Conversação'
tf = body_shape.text_frame
tf.text = 'A IA consegue conversar com pessoas de forma natural, como uma IA comum.'
p = tf.add_paragraph()
p.text = 'Ela responde a perguntas, fornece explicações e interage de maneira amigável.'
p.level = 0

# Slide 5: Benefícios
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = 'Benefícios'
tf = body_shape.text_frame
tf.text = 'Aumento da produtividade no desenvolvimento de software.'
p = tf.add_paragraph()
p.text = 'Redução de erros humanos no código.'
p.level = 0
p = tf.add_paragraph()
p.text = 'Correção rápida de problemas existentes.'
p.level = 0

# Slide 6: Otimização do Sistema
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = 'Otimização do Sistema'
tf = body_shape.text_frame
tf.text = 'A IA otimiza o sistema automaticamente a cada hora, pesquisando e executando comandos.'
p = tf.add_paragraph()
p.text = 'Inclui otimização de disco, memória, limpeza de cache e verificação antivírus.'
p.level = 0
p = tf.add_paragraph()
p.text = 'Ativa permissões necessárias (sudo) para tarefas administrativas.'
p.level = 0

# Slide 7: Conclusão
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = 'Conclusão'
tf = body_shape.text_frame
tf.text = 'A IA IARA-OS-ENGINE representa o futuro da programação assistida por IA.'
p = tf.add_paragraph()
p.text = 'Ela permite criar e manter aplicações de alta qualidade com eficiência.'
p.level = 0

# Salvar a apresentação
prs.save('apresentacao_ia.pptx')